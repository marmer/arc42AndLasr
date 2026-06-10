#!/usr/bin/env python3
"""Build the reveal.js presentation from the original talk material.

Pipeline:
  1. Open the rendered PDF of the original talk (one page per visible slide).
  2. Strip company branding:
       - the company logo image (appears on every page),
       - the company tagline text,
       - the company e-mail address (replaced with the private one),
       - the vCard QR code (regenerated without company data).
  3. Export every page as a vector SVG (text as paths, so no font
     dependencies) to docs/img/slide-NN.svg.
  4. Pull slide titles and speaker notes from the PPTX, skipping hidden
     slides so the numbering matches the PDF pages.
  5. Generate docs/index.html with one <section> per slide.

Usage:
    pip install -r scripts/requirements.txt
    python3 scripts/build_presentation.py
"""

import html
import io
import re
import sys
from pathlib import Path

import fitz  # PyMuPDF
import qrcode
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PDF = ROOT / "arc42AndLasr_talk - envite_original_rendered.pdf"
PPTX = ROOT / "arc42AndLasr_talk - envite_original.pptx"
DOCS = ROOT / "docs"
IMG = DOCS / "img"
KARLA = Path(__file__).resolve().parent / "assets" / "Karla.ttf"

BANNED = re.compile(r"(?i)novatec|envite|it-tage|software quality days|decompiled")

TAGLINE = "Pioneering IT Sustainability"
OLD_EMAIL = "Mariano.Mertinat@envite.de"
NEW_EMAIL = "mariano.mertinat@gmail.com"
EMAIL_COLOR = (0x1B / 255, 0x1B / 255, 0x1B / 255)
CARD_FILL = (108 / 255, 202 / 255, 177 / 255)  # teal contact card on the last slide

VCARD = (
    "BEGIN:VCARD\n"
    "VERSION:2.1\n"
    "N;CHARSET=UTF-8:Mertinat;Mariano\n"
    "FN;CHARSET=UTF-8:Mariano Mertinat\n"
    "TEL;HOME;VOICE:+491707095722\n"
    "TITLE;CHARSET=UTF-8:green IT-consultant / software-architect / "
    "solution-architect / trainer\n"
    f"EMAIL:{NEW_EMAIL}\n"
    "URL:https://arc42andlasr.marmer.online\n"
    "END:VCARD"
)


def find_logo_xref(doc):
    """The company logo is the image that is placed on (nearly) every page."""
    usage = {}
    for page in doc:
        for img in page.get_images(full=True):
            usage[img[0]] = usage.get(img[0], 0) + 1
    xref, count = max(usage.items(), key=lambda kv: kv[1])
    if count < len(doc) * 0.8:
        sys.exit(f"no image appears on most pages (best: xref {xref} on {count} pages)")
    return xref


def find_vcard_qr(page):
    """The vCard QR code is the square image on the contact card of the last page."""
    for img in page.get_images(full=True):
        bbox = page.get_image_bbox(img)
        if abs(bbox.width - bbox.height) < 5 and 100 < bbox.width < 200 and bbox.x0 > 300:
            return img[0], bbox
    sys.exit("vCard QR code not found on the last page")


def clean_pdf(doc):
    logo_xref = find_logo_xref(doc)
    doc[0].delete_image(logo_xref)  # makes the image transparent document-wide

    for page in doc:
        rects = page.search_for(TAGLINE)
        if rects:
            for r in rects:
                page.add_redact_annot(r, fill=(1, 1, 1))
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

    last = doc[-1]
    email_rects = last.search_for(OLD_EMAIL)
    if not email_rects:
        sys.exit("company e-mail not found on the last page")
    origin = fitz.Point(email_rects[0].x0, email_rects[0].y1 - 3.8)
    for r in email_rects:
        last.add_redact_annot(r, fill=CARD_FILL)
    last.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
    last.insert_text(origin, NEW_EMAIL, fontsize=15.0, fontname="Karla",
                     fontfile=str(KARLA), color=EMAIL_COLOR)

    qr_xref, _ = find_vcard_qr(last)
    qr_img = qrcode.make(VCARD, border=1).get_image()
    buf = io.BytesIO()
    qr_img.save(buf, format="PNG")
    last.replace_image(qr_xref, stream=buf.getvalue())

    for page in doc:
        text = page.get_text()
        if BANNED.search(text):
            sys.exit(f"banned reference still present on page {page.number + 1}: "
                     f"{BANNED.search(text).group(0)}")


def visible_slides(prs):
    return [s for s in prs.slides if s._element.get("show") != "0"]


def slide_title(slide):
    if slide.shapes.title is not None and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip().splitlines()[0]
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip().splitlines()[0]
    return ""


def slide_notes(slide):
    if not slide.has_notes_slide:
        return ""
    text = slide.notes_slide.notes_text_frame.text.strip()
    if BANNED.search(text):
        sys.exit(f"banned reference in speaker notes: {text[:80]}")
    return text


def notes_html(text):
    if not text:
        return ""
    paragraphs = [html.escape(p.strip()) for p in text.split("\n") if p.strip()]
    return "\n                ".join(f"<p>{p}</p>" for p in paragraphs)


SECTION_TEMPLATE = """\
        <!-- Slide {num}: {title} -->
        <section aria-label="{title_attr}">
            <img {src_attr}="./img/slide-{num:02d}.svg" class="slide-img" alt="{title_attr}">
            <aside aria-label="speaker notes" class="notes">
                {notes}
            </aside>
        </section>
"""

PAGE_TEMPLATE = """\
<!doctype html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no">

    <title>arc42 and LASR — Understand and improve systems in a playful way</title>

    <link rel="stylesheet" href="dist/reset.css">
    <link rel="stylesheet" href="dist/reveal.css">
    <link rel="stylesheet" href="dist/theme/white.css">
    <link rel="stylesheet" href="css/custom.css">
</head>
<body>
<div class="reveal">
    <div class="slides">

{sections}
    </div>
</div>

<script src="dist/reveal.js"></script>
<script src="dist/plugin/notes.js"></script>
<script src="dist/plugin/markdown.js"></script>
<script src="dist/plugin/search.js"></script>
<script src="dist/plugin/zoom.js"></script>
<script src="dist/plugin/highlight.js"></script>
<script>
    Reveal.initialize({{
        width: 960,
        height: 540,
        margin: 0,
        controls: false,
        progress: false,
        hash: true,
        slideNumber: "c/t",
        history: true,
        mouseWheel: true,
        transition: 'fade',
        navigationMode: "linear",
        plugins: [ RevealMarkdown, RevealHighlight, RevealNotes, RevealSearch, RevealZoom ]
    }});
</script>
</body>
</html>
"""


def build():
    doc = fitz.open(PDF)
    prs = Presentation(PPTX)
    slides = visible_slides(prs)
    if len(slides) != len(doc):
        sys.exit(f"slide count mismatch: {len(slides)} visible PPTX slides, "
                 f"{len(doc)} PDF pages")

    clean_pdf(doc)

    IMG.mkdir(parents=True, exist_ok=True)
    sections = []
    for i, (page, slide) in enumerate(zip(doc, slides), 1):
        svg = page.get_svg_image(text_as_path=True)
        (IMG / f"slide-{i:02d}.svg").write_text(svg)
        title = slide_title(slide) or f"Slide {i}"
        sections.append(SECTION_TEMPLATE.format(
            num=i,
            title=title,
            title_attr=html.escape(title, quote=True),
            # the first slide loads eagerly so it shows up immediately,
            # the rest lazily via reveal.js data-src
            src_attr="src" if i == 1 else "data-src",
            notes=notes_html(slide_notes(slide)) or "",
        ))
        print(f"slide {i:02d}: {title}")

    (DOCS / "index.html").write_text(PAGE_TEMPLATE.format(sections="\n".join(sections)))
    print(f"\nwrote {len(sections)} slides to {DOCS / 'index.html'}")


if __name__ == "__main__":
    build()
